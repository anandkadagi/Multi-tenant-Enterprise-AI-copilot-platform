from pathlib import Path
import filetype
import pandas as pd

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
def load(file_path):

    file_path = Path(file_path)
    doc_type = detect_type(
        file_path
    )
    pages = extract_text(
        file_path,
        doc_type
    )
    return {
        "filename": file_path.name,
        "path": str(file_path),
        "document_type": doc_type,
        "pages": pages
    }
def detect_type(file_path):

    try:
        kind = filetype.guess(
        str(file_path)
        )

        if kind:
            return kind.mime

       
    except:
        mapping = {
            ".pdf": "application/pdf",
            ".docx": "docx",
            ".txt": "text",
            ".csv": "csv",
            ".xlsx": "excel",
            ".pptx": "pptx"
        }
        return mapping.get(
            file_path.suffix.lower(),
            "unknown"
        )
def extract_text(
        file_path,
        mime
    ):

    if "pdf" in mime:
        return parse_pdf(
            file_path
        )
    elif "word" in mime:
        return parse_docx(
            file_path
        )
    elif "ppt" in mime:
        return parse_pptx(
            file_path
        )
    elif "spreadsheet" in mime:
        return parse_excel(
            file_path
        )
    elif "csv" in mime:
        return parse_csv(
            file_path
        )
    elif "text" in mime:
        return parse_txt(
            file_path
        )
    return []
# -------------------
    # PDF
    # -------------------
def parse_pdf(
    path
):
    reader = PdfReader(
        str(path)
    )
    pages = []
    for idx, page in enumerate(
        reader.pages,
        start=1
    ):
        pages.append(
            {
                "page_number": idx,
                "text":
                page.extract_text()
                or ""
            }
        )
    return pages
# -------------------
# DOCX
# -------------------
def parse_docx(
    path
):
    doc = Document(
        str(path)
    )
    return [
        {
            "page_number": 1,
            "text": "\n".join(
                p.text
                for p in doc.paragraphs
            )
        }
    ]
# -------------------
# TXT
# -------------------
def parse_txt(
    path
):
    return [
        {
            "page_number": 1,
            "text":
            path.read_text(
                encoding="utf-8",
                errors="ignore"
            )
        }
    ]
# -------------------
# CSV
# -------------------
def parse_csv(
    path
):
    df = pd.read_csv(
        path
    )
    return [
        {
            "page_number": 1,
            "text":
            df.to_string(
                index=False
            )
        }
    ]
# -------------------
# Excel
# -------------------
def parse_excel(
    path
):
    sheets = pd.read_excel(
        path,
        sheet_name=None
    )
    result = []
    for sheet_no, (
        name,
        df
    ) in enumerate(
        sheets.items(),
        start=1
    ):
        result.append(
            {
                "page_number": sheet_no,
                "sheet": name,
                "text":
                df.to_string(
                    index=False
                )
            }
        )
    return result
# -------------------
# PPT
# -------------------
def parse_pptx(
    path
):
    prs = Presentation(
        str(path)
    )
    slides = []
    for i, slide in enumerate(
        prs.slides,
        start=1
    ):
        txt = []
        for shape in slide.shapes:
            if hasattr(
                shape,
                "text"
            ):
                txt.append(
                    shape.text
                )
        slides.append(
            {
                "page_number": i,
                "text":
                "\n".join(
                    txt
                )
            }
        )
    return slides